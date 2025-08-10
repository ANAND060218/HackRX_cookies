# app/cache_manager.py
import os
import pickle
import hashlib
from pathlib import Path


CACHE_DIR = Path(os.getenv("CACHE_DIR", "./vector_cache"))
CACHE_DIR.mkdir(parents=True, exist_ok=True)

MAPPING_FILE = CACHE_DIR / "file_mapping.pkl"  # keeping name as is

if MAPPING_FILE.exists():
    with MAPPING_FILE.open("rb") as f:
        url_mapping = pickle.load(f)
else:
    url_mapping = {}
#-----------------------------------#---------------#

def _compute_file_text_hash(file_url: str, snippet_len: int = 16) -> str:
    """
    Quickly extract only enough text to generate a hash.
    For PDFs, DOCX, PPTX, XLSX, read minimal pages/slides/rows.
    """
    import requests, mimetypes
    from tempfile import NamedTemporaryFile
    from pathlib import Path

    # Download file to temp
    head = requests.head(file_url, allow_redirects=True, timeout=10)
    ext = file_url.split('?')[0].split('.')[-1].lower()
    if not ext:
        ext = (mimetypes.guess_extension(head.headers.get("Content-Type", "")) or "bin").lstrip('.')

    r = requests.get(file_url, timeout=30)
    r.raise_for_status()
    with NamedTemporaryFile(delete=False, suffix=f".{ext}") as f:
        f.write(r.content)
        f.flush()
        tmp_path = Path(f.name)

    text = ""
    try:
        if ext == "pdf":
            import fitz
            doc = fitz.open(str(tmp_path))
            max_pages = min(len(doc), 2)  # only first 2 pages
            text = "\n".join(doc[i].get_text() for i in range(max_pages))

        elif ext == "docx":
            import docx
            doc = docx.Document(str(tmp_path))
            text = "\n".join(p.text for p in doc.paragraphs[:5])  # first 5 paragraphs

        elif ext == "pptx":
            from pptx import Presentation
            prs = Presentation(str(tmp_path))
            slides = []
            for slide in prs.slides[:2]:  # first 2 slides
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    slides.append("\n".join(slide_text))
            text = "\n".join(slides)

        elif ext in ("xlsx", "xls"):
            import openpyxl
            wb = openpyxl.load_workbook(str(tmp_path), data_only=True)
            rows = []
            for sheet in wb.worksheets[:1]:  # first sheet only
                for i, row in enumerate(sheet.iter_rows(values_only=True)):
                    if i >= 5:  # first 5 rows
                        break
                    rows.append(" ".join(str(c) if c is not None else "" for c in row))
            text = "\n".join(rows)

        else:
            text = tmp_path.read_text(errors="ignore")[:snippet_len]

    except Exception:
        return None

    # Normalize and hash
    normalized = text.strip().replace("\n", " ").replace("\r", " ").lower()
    snippet = normalized[:snippet_len]
    import hashlib
    return hashlib.sha256(snippet.encode("utf-8")).hexdigest()

def load_vector_store_if_exists(url: str):
    filehash = _compute_file_text_hash(url)
    if not filehash:
        return None
    path = url_mapping.get(filehash)
    if path and Path(path).exists():
        with open(path, "rb") as f:
            return pickle.load(f)
    return None
#-----------------------------------#---------------#
def save_vector_store(db, url: str):
    """
    Save picklable `db` under a numbered folder keyed by file content hash.
    """
    filehash = _compute_file_text_hash(url)
    if not filehash:
        return

    if filehash in url_mapping:
        path = Path(url_mapping[filehash])
    else:
        next_id = len(url_mapping) + 1
        dir_path = CACHE_DIR / f"vector_store_{next_id}"
        dir_path.mkdir(parents=True, exist_ok=True)
        path = dir_path / "db.pkl"
        url_mapping[filehash] = str(path)
        with MAPPING_FILE.open("wb") as f:
            pickle.dump(url_mapping, f)

    with open(path, "wb") as f:
        pickle.dump(db, f)
#-----------------------------------#---------------#